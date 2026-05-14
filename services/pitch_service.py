# """
# services/pitch_service.py
# Generates a BizGenius investor pitch deck (.pptx) using python-pptx.
# No Node.js / npm required.
# """

# import os
# import tempfile
# import json
# from datetime import datetime
# from pptx import Presentation
# from pptx.util import Inches, Pt, Emu
# from pptx.dml.color import RGBColor
# from pptx.enum.text import PP_ALIGN


# # ── Palette ──────────────────────────────────────────────────
# BG      = RGBColor(0x0d, 0x0d, 0x24)
# PURPLE  = RGBColor(0x7c, 0x6b, 0xff)
# WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
# GRAY    = RGBColor(0x94, 0xa3, 0xb8)
# SURFACE = RGBColor(0x11, 0x11, 0x30)

# _ACCENT_MAP = {
#     "Success":  (RGBColor(0x34, 0xd3, 0x99), RGBColor(0x06, 0x5f, 0x46)),
#     "Failure":  (RGBColor(0xf8, 0x71, 0x71), RGBColor(0x7f, 0x1d, 0x1d)),
#     "Uncertain":(RGBColor(0xfb, 0xbf, 0x24), RGBColor(0x78, 0x35, 0x0f)),
# }


# # ── Low-level helpers ─────────────────────────────────────────

# def _rgb(r, g, b):
#     return RGBColor(r, g, b)

# def _fill(shape, color: RGBColor):
#     shape.fill.solid()
#     shape.fill.fore_color.rgb = color

# def _no_line(shape):
#     shape.line.fill.background()

# def _box(slide, x, y, w, h, fill_color: RGBColor, line_color=None, line_width_pt=0):
#     from pptx.util import Pt
#     shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
#     _fill(shape, fill_color)
#     if line_color:
#         shape.line.color.rgb = line_color
#         shape.line.width = Pt(line_width_pt)
#     else:
#         _no_line(shape)
#     return shape

# def _text(slide, text, x, y, w, h,
#           font_size=11, color=WHITE, bold=False, italic=False,
#           align=PP_ALIGN.LEFT, wrap=True):
#     txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
#     txBox.word_wrap = wrap
#     tf = txBox.text_frame
#     tf.word_wrap = wrap
#     p = tf.paragraphs[0]
#     p.alignment = align
#     run = p.add_run()
#     run.text = str(text)
#     run.font.size = Pt(font_size)
#     run.font.color.rgb = color
#     run.font.bold = bold
#     run.font.italic = italic
#     run.font.name = "Calibri"
#     return txBox

# def _multiline_text(slide, lines, x, y, w, h,
#                     font_size=10, color=GRAY, bold=False):
#     """Add a textbox with multiple lines as separate paragraphs."""
#     txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
#     txBox.word_wrap = True
#     tf = txBox.text_frame
#     tf.word_wrap = True
#     for i, line in enumerate(lines):
#         p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
#         p.alignment = PP_ALIGN.LEFT
#         run = p.add_run()
#         run.text = str(line)
#         run.font.size = Pt(font_size)
#         run.font.color.rgb = color
#         run.font.bold = bold
#         run.font.name = "Calibri"
#     return txBox

# def _bg(slide):
#     """Fill entire slide background."""
#     _box(slide, 0, 0, 10, 5.625, BG)

# def _accent_bar(slide, color: RGBColor):
#     """Left edge accent bar."""
#     _box(slide, 0, 0, 0.08, 5.625, color)

# def _top_bar(slide, color: RGBColor):
#     _box(slide, 0, 0, 10, 0.07, color)

# def _section_title(slide, title: str, accent: RGBColor):
#     _box(slide, 0, 0, 10, 0.07, accent)
#     _text(slide, title, 0.4, 0.15, 9.2, 0.55,
#           font_size=24, color=WHITE, bold=True)


# # ── Slide builders ────────────────────────────────────────────

# def _slide1_title(prs, domain, description, classification, sp, accent):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
#     _bg(slide)
#     _accent_bar(slide, PURPLE)

#     _text(slide, domain.upper(), 0.4, 0.45, 9, 0.4,
#           font_size=11, color=accent, bold=True)
#     _text(slide, "Startup Pitch Deck", 0.4, 0.9, 9, 1.1,
#           font_size=40, color=WHITE, bold=True)
#     _text(slide, description[:90], 0.4, 2.1, 8, 0.7,
#           font_size=14, color=GRAY, italic=True)

#     # Classification badge
#     _box(slide, 0.4, 3.1, 2.2, 0.42, PURPLE)
#     _text(slide, f"{classification}  •  {sp:.0f}% Success",
#           0.42, 3.13, 2.16, 0.36,
#           font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

#     _text(slide, datetime.now().strftime("%B %Y"),
#           0.4, 5.1, 4, 0.35, font_size=10, color=GRAY)


# def _slide2_problem(prs, ui, description, domain):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide)
#     _section_title(slide, "The Problem & Opportunity", PURPLE)

#     # Left panel — problem
#     _box(slide, 0.4, 0.85, 4.3, 3.8, SURFACE)
#     _text(slide, "PROBLEM", 0.55, 0.95, 4.0, 0.35,
#           font_size=10, color=RGBColor(0x34,0xd3,0x99), bold=True)
#     _text(slide, description[:220], 0.55, 1.38, 4.0, 3.0,
#           font_size=10, color=GRAY, wrap=True)

#     # Right top — domain
#     _box(slide, 5.1, 0.85, 4.5, 1.7, SURFACE)
#     _text(slide, "DOMAIN", 5.25, 0.95, 4.0, 0.35,
#           font_size=10, color=RGBColor(0x34,0xd3,0x99), bold=True)
#     _text(slide, domain, 5.25, 1.38, 4.0, 0.8,
#           font_size=22, color=WHITE, bold=True)

#     # Right bottom — key stats
#     _box(slide, 5.1, 2.65, 4.5, 2.0, SURFACE)
#     _text(slide, "KEY STATS", 5.25, 2.75, 4.0, 0.35,
#           font_size=10, color=RGBColor(0x34,0xd3,0x99), bold=True)
#     stats_lines = [
#         f"Age: {ui.get('company_age')} yrs   Founders: {ui.get('founder_count')}",
#         f"Employees: {ui.get('employees')}   Investors: {ui.get('investor_count')}",
#         f"Funding Rounds: {ui.get('funding_rounds')}",
#     ]
#     _multiline_text(slide, stats_lines, 5.25, 3.18, 4.2, 1.3,
#                     font_size=10, color=GRAY)


# def _slide3_ml(prs, ml, accent):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide)
#     _section_title(slide, "ML Prediction Results", PURPLE)

#     classification = ml.get("classification", "—")
#     sp   = ml.get("success_probability", 0) * 100
#     risk = ml.get("risk_level", "—")
#     pred = ml.get("predicted_funding_usd", 0)

#     cards = [
#         ("Classification",      classification,         accent),
#         ("Success Probability", f"{sp:.1f}%",            PURPLE),
#         ("Risk Level",          risk,                   accent),
#         ("Predicted Funding",   f"${pred/1e6:.2f}M",    RGBColor(0x81,0x8c,0xf8)),
#     ]

#     positions = [(0.4, 1.0), (5.1, 1.0), (0.4, 2.95), (5.1, 2.95)]
#     for (x, y), (label, value, col) in zip(positions, cards):
#         _box(slide, x, y, 4.3, 1.6, SURFACE, line_color=col, line_width_pt=2)
#         _text(slide, label.upper(), x+0.15, y+0.15, 4.0, 0.35,
#               font_size=9, color=GRAY)
#         _text(slide, value, x+0.15, y+0.55, 4.0, 0.75,
#               font_size=26, color=col, bold=True)


# def _slide4_competitors(prs, comp_names, accent):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide)
#     _section_title(slide, "Competitor Landscape", PURPLE)

#     if not comp_names:
#         _text(slide, "No competitor data available.", 0.4, 2.0, 9.2, 0.5,
#               font_size=13, color=GRAY, italic=True)
#         return

#     for i, name in enumerate(comp_names[:6]):
#         y = 0.98 + i * 0.68
#         _box(slide, 0.4, y, 9.2, 0.55, SURFACE)
#         _box(slide, 0.4, y, 0.07, 0.55, PURPLE)
#         _text(slide, f"{i+1}.  {name}", 0.6, y+0.1, 8.8, 0.38,
#               font_size=10, color=WHITE)


# def _slide5_team(prs, hierarchy, accent):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide)
#     _section_title(slide, "Team & Organization", PURPLE)

#     ceo   = hierarchy.get("ceo_title", "CEO")
#     total = hierarchy.get("total_employees", "—")
#     _text(slide, f"CEO: {ceo}   •   Total: {total} employees",
#           0.4, 0.82, 9, 0.4, font_size=12, color=accent, bold=True)

#     depts = hierarchy.get("departments", [])
#     for i, d in enumerate(depts[:5]):
#         y   = 1.3 + i * 0.63
#         col = SURFACE if i % 2 == 0 else RGBColor(0x0d, 0x1a, 0x35)
#         roles = ", ".join(d.get("roles", [])[:2])
#         line  = f"{d.get('name','')}:  {d.get('headcount','')} people  —  {roles}"
#         _box(slide, 0.4, y, 9.2, 0.55, col)
#         _text(slide, line, 0.6, y+0.1, 8.8, 0.38, font_size=10, color=WHITE)


# def _slide6_risks(prs, risks, accent):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide)
#     _section_title(slide, "Risk Management", PURPLE)

#     for i, risk in enumerate(risks[:6]):
#         y = 0.98 + i * 0.68
#         _box(slide, 0.4, y, 9.2, 0.55, SURFACE)
#         _box(slide, 0.4, y, 0.07, 0.55, accent)
#         _text(slide, f"⚠   {risk}", 0.6, y+0.1, 8.8, 0.38,
#               font_size=10, color=WHITE)


# def _slide7_action_plan(prs, plan_lines, accent):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide)
#     _section_title(slide, "30-Day Action Plan", PURPLE)

#     for i, task in enumerate(plan_lines[:6]):
#         x = 0.4 if i < 3 else 5.1
#         y = 1.0 + (i % 3) * 1.35
#         _box(slide, x, y, 4.5, 1.2, SURFACE)
#         _text(slide, f"DAY {(i+1)*5}", x+0.15, y+0.08, 4.2, 0.3,
#               font_size=9, color=accent, bold=True)
#         _text(slide, task, x+0.15, y+0.42, 4.2, 0.68,
#               font_size=10, color=WHITE, wrap=True)


# def _slide8_roadmap(prs, hire_lines, pred_fund, accent):
#     slide = prs.slides.add_slide(prs.slide_layouts[6])
#     _bg(slide)
#     _top_bar(slide, PURPLE)
#     _section_title(slide, "Future Roadmap & The Ask", PURPLE)

#     # Key hires
#     _box(slide, 0.4, 0.95, 4.5, 2.2, SURFACE)
#     _text(slide, "KEY HIRES NEEDED", 0.55, 1.05, 4.2, 0.38,
#           font_size=11, color=accent, bold=True)
#     _multiline_text(slide, [f"• {h}" for h in hire_lines[:3]],
#                     0.55, 1.5, 4.2, 1.5, font_size=9, color=GRAY)

#     # Funding ask
#     _box(slide, 5.1, 0.95, 4.5, 2.2, SURFACE)
#     _text(slide, "FUNDING ASK", 5.25, 1.05, 4.2, 0.38,
#           font_size=11, color=accent, bold=True)
#     _text(slide, f"${pred_fund/1e6:.2f}M", 5.25, 1.5, 4.2, 0.8,
#           font_size=34, color=WHITE, bold=True)
#     _text(slide, "Predicted Next Round", 5.25, 2.38, 4.2, 0.4,
#           font_size=10, color=GRAY)

#     # Milestones
#     _box(slide, 0.4, 3.3, 9.2, 1.8, SURFACE)
#     _text(slide, "GROWTH MILESTONES", 0.55, 3.4, 8.8, 0.38,
#           font_size=11, color=PURPLE, bold=True)
#     milestones = [
#         "Q1: MVP launch & first 100 customers",
#         "Q2: Product-market fit validation",
#         "Q3: Scale to $1M ARR",
#         "Q4: Series A raise",
#     ]
#     _multiline_text(slide, milestones, 0.55, 3.85, 8.8, 1.0,
#                     font_size=11, color=GRAY)


# # ── Public entry point ────────────────────────────────────────

# def generate_pitch_deck(data: dict) -> str:
#     """
#     Build the pitch deck and return the path to the saved .pptx file.
#     Caller is responsible for deleting the file after streaming.
#     """
#     ui          = data["user_input"]
#     ml          = data["ml_results"]
#     analysis    = data.get("analysis", "")
#     competitors = data.get("competitors", [])
#     hierarchy   = data.get("hierarchy", {})
#     risks       = data.get("probable_risks", [])

#     classification = ml.get("classification", "Uncertain")
#     sp             = ml.get("success_probability", 0) * 100
#     pred_fund      = ml.get("predicted_funding_usd", 0)
#     domain         = ui.get("domain", "")
#     description    = ui.get("description", "")

#     accent, _ = _ACCENT_MAP.get(classification, _ACCENT_MAP["Uncertain"])

#     # ── Parse 30-day plan from LLM analysis ──
#     plan_lines = []
#     in_plan = False
#     for line in analysis.split("\n"):
#         if "30-day" in line.lower() or "30 day" in line.lower():
#             in_plan = True
#         if in_plan and line.strip().startswith(("-", "•", "*")):
#             plan_lines.append(line.strip().lstrip("-•* "))
#         if in_plan and len(plan_lines) >= 6:
#             break
#     if not plan_lines:
#         plan_lines = [
#             "Validate core assumptions with 10 customer interviews",
#             "Define MVP feature set and ship v0.1",
#             "Set up analytics and growth tracking",
#             "Identify top 3 acquisition channels",
#             "Establish weekly team cadence and OKRs",
#             "Prepare pitch materials for next funding round",
#         ]

#     # ── Competitor display names ──
#     comp_names = []
#     for c in competitors[:6]:
#         meta     = c.get("metadata", {})
#         industry = meta.get("industry", "")
#         funding  = meta.get("funding", 0)
#         doc      = c.get("document", "")
#         if funding:
#             comp_names.append(f"{industry} — ${funding/1e6:.1f}M funded")
#         else:
#             comp_names.append(doc[:70])

#     # ── Hire display lines ──
#     next_hires = hierarchy.get("recommended_next_hires", [])
#     hire_lines = [
#         f"{h['role']} ({h['priority']}): {h['reason']}"
#         for h in next_hires[:3]
#     ]

#     # ── Build deck ──
#     prs = Presentation()
#     prs.slide_width  = Inches(10)
#     prs.slide_height = Inches(5.625)

#     _slide1_title(prs, domain, description, classification, sp, accent)
#     _slide2_problem(prs, ui, description, domain)
#     _slide3_ml(prs, ml, accent)
#     _slide4_competitors(prs, comp_names, accent)
#     _slide5_team(prs, hierarchy, accent)
#     _slide6_risks(prs, risks, accent)
#     _slide7_action_plan(prs, plan_lines, accent)
#     _slide8_roadmap(prs, hire_lines, pred_fund, accent)

#     # ── Save to temp file ──
#     tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
#     prs.save(tmp.name)
#     tmp.close()
#     return tmp.name

"""
services/pitch_service.py  —  BizGenius Investor Pitch Deck
Mirrors every section of the PDF report:
  1. Cover / Startup Snapshot
  2. ML Prediction Results
  3. Risk Factors
  4. Competitor Analysis
  5. Team Hierarchy & Org Structure
  6. Hiring Guide & Role Profiles
  7. Strategic Analysis & 30-Day Action Plan
  8. Closing / Footer
"""

import os
import tempfile
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Palette (matches report_service.py) ──────────────────────
BG_DARK  = RGBColor(0x0d, 0x0d, 0x24)
BG_MID   = RGBColor(0x11, 0x11, 0x30)
BG_CARD  = RGBColor(0x16, 0x16, 0x3a)
PURPLE   = RGBColor(0x7c, 0x6b, 0xff)
PURPLE_D = RGBColor(0x4c, 0x3f, 0xcc)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x94, 0xa3, 0xb8)
GRAY_L   = RGBColor(0xcb, 0xd5, 0xe1)
GREEN    = RGBColor(0x34, 0xd3, 0x99)
RED      = RGBColor(0xf8, 0x71, 0x71)
YELLOW   = RGBColor(0xfb, 0xbf, 0x24)
ORANGE   = RGBColor(0xfb, 0x92, 0x3c)
BLACK    = RGBColor(0x1e, 0x1b, 0x4b)

_ACCENT_MAP = {
    "Success":   GREEN,
    "Failure":   RED,
    "Uncertain": YELLOW,
}

W  = 10.0   # slide width  inches
H  = 5.625  # slide height inches


# ── Low-level drawing helpers ─────────────────────────────────

def _fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def _no_line(shape):
    shape.line.fill.background()

def _rect(slide, x, y, w, h, fill: RGBColor, border: RGBColor = None, border_pt=1.5):
    from pptx.util import Pt as _Pt
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(s, fill)
    if border:
        s.line.color.rgb = border
        s.line.width = _Pt(border_pt)
    else:
        _no_line(s)
    return s

def _txt(slide, text, x, y, w, h,
         size=11, color=WHITE, bold=False, italic=False,
         align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return tb

def _multiline(slide, lines, x, y, w, h,
               size=10, color=GRAY, bold=False, font="Calibri",
               align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = str(line)
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
    return tb

def _bg(slide):
    _rect(slide, 0, 0, W, H, BG_DARK)

def _slide_header(slide, title: str, slide_num: int, total: int):
    """Top strip with slide title + number."""
    _rect(slide, 0, 0, W, 0.55, PURPLE)
    _txt(slide, title, 0.3, 0.07, 8.5, 0.42, size=18, bold=True,
         color=WHITE, align=PP_ALIGN.LEFT)
    _txt(slide, f"{slide_num}/{total}", 9.1, 0.1, 0.7, 0.35,
         size=10, color=WHITE, align=PP_ALIGN.RIGHT)

def _label(slide, text, x, y, color=GREEN):
    """Small ALL-CAPS section label."""
    _txt(slide, text.upper(), x, y, 4.5, 0.3, size=8, color=color, bold=True)

def _card(slide, x, y, w, h, accent: RGBColor = None):
    """Dark card with optional left accent stripe."""
    _rect(slide, x, y, w, h, BG_CARD)
    if accent:
        _rect(slide, x, y, 0.06, h, accent)

def _stat_card(slide, x, y, w, h, label, value, accent: RGBColor):
    """Big-number stat card."""
    _card(slide, x, y, w, h, accent)
    _txt(slide, label.upper(), x+0.18, y+0.12, w-0.25, 0.3,
         size=8, color=GRAY, bold=True)
    _txt(slide, str(value), x+0.18, y+0.45, w-0.25, h-0.55,
         size=22, color=accent, bold=True)

def _trunc(text, n=120):
    t = str(text or "")
    return t[:n] + "…" if len(t) > n else t

def _clean(text):
    """Strip markdown and clean for slides."""
    import re
    t = str(text or "")
    t = re.sub(r'\*\*(.+?)\*\*', r'\1', t)
    t = re.sub(r'\*(.+?)\*',   r'\1', t)
    t = re.sub(r'#+\s*',       '',    t)
    return t.strip()


# ════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ════════════════════════════════════════════════════════════════

TOTAL_SLIDES = 10   # used for slide counter

# ── Slide 1 · Cover ──────────────────────────────────────────
def _slide_cover(prs, ui, ml):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    domain = ui.get("domain", "")
    desc   = ui.get("description", "")
    cls    = ml.get("classification", "Uncertain")
    sp     = ml.get("success_probability", 0) * 100
    accent = _ACCENT_MAP.get(cls, YELLOW)

    # Left decorative stripe
    _rect(slide, 0, 0, 0.12, H, PURPLE)

    # Large title
    _txt(slide, "BizGenius", 0.5, 0.5, 7, 0.7,
         size=38, bold=True, color=WHITE, font="Calibri")
    _txt(slide, "Startup Intelligence Report", 0.5, 1.2, 7, 0.5,
         size=18, color=GRAY, font="Calibri")

    # Domain badge
    _rect(slide, 0.5, 1.85, 3.2, 0.42, PURPLE_D)
    _txt(slide, domain.upper(), 0.55, 1.9, 3.1, 0.32,
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Description
    _txt(slide, _trunc(desc, 160), 0.5, 2.5, 6.2, 0.9,
         size=11, color=GRAY_L, italic=True, wrap=True)

    # Classification badge bottom-left
    _rect(slide, 0.5, 4.6, 3.5, 0.7, BG_CARD)
    _rect(slide, 0.5, 4.6, 0.06, 0.7, accent)
    _txt(slide, f"{cls}  •  {sp:.0f}% Success Probability",
         0.65, 4.72, 3.2, 0.35, size=11, color=accent, bold=True)

    # Date bottom-right
    _txt(slide, datetime.now().strftime("%B %Y"),
         7.5, 5.1, 2.3, 0.35, size=10, color=GRAY, align=PP_ALIGN.RIGHT)

    # Right panel — quick key stats
    _rect(slide, 7.0, 0, 3.0, H, BG_MID)
    _txt(slide, "KEY FACTS", 7.2, 0.3, 2.6, 0.35,
         size=10, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    stats = [
        ("Company Age",    f"{ui.get('company_age','—')} yrs"),
        ("Founders",       str(ui.get("founder_count","—"))),
        ("Employees",      str(ui.get("employees","—"))),
        ("Funding Rounds", str(ui.get("funding_rounds","—"))),
        ("Avg per Round",  f"${ui.get('funding_per_round',0):,.0f}"),
        ("Investors",      str(ui.get("investor_count","—"))),
    ]
    for i, (lbl, val) in enumerate(stats):
        y = 0.75 + i * 0.72
        _rect(slide, 7.15, y, 2.6, 0.6, BG_CARD)
        _txt(slide, lbl, 7.25, y+0.04, 2.4, 0.25, size=8, color=GRAY)
        _txt(slide, val, 7.25, y+0.27, 2.4, 0.28, size=13, bold=True, color=WHITE)


# ── Slide 2 · Startup Snapshot ────────────────────────────────
def _slide_snapshot(prs, ui, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "1. Startup Snapshot", n, TOTAL_SLIDES)

    rows = [
        ("Domain",             ui.get("domain","—")),
        ("Description",        _trunc(ui.get("description","—"), 180)),
        ("Company Age",        f"{ui.get('company_age','—')} years"),
        ("Founders",           str(ui.get("founder_count","—"))),
        ("Employees",          str(ui.get("employees","—"))),
        ("Funding Rounds",     str(ui.get("funding_rounds","—"))),
        ("Avg Funding / Round",f"${ui.get('funding_per_round',0):,.0f}"),
        ("Investors",          str(ui.get("investor_count","—"))),
    ]

    col_colors = [PURPLE, BG_CARD]
    for i, (field, val) in enumerate(rows):
        y = 0.65 + i * 0.575
        _rect(slide, 0.3,  y, 2.8, 0.52, col_colors[i % 2])
        _rect(slide, 3.15, y, 6.5, 0.52, BG_CARD)
        _txt(slide, field, 0.42, y+0.1, 2.6, 0.35, size=10, bold=True, color=WHITE)
        _txt(slide, val,   3.28, y+0.1, 6.2, 0.35, size=10, color=GRAY_L, wrap=True)


# ── Slide 3 · ML Prediction Results ──────────────────────────
def _slide_ml(prs, ml, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "2. ML Prediction Results", n, TOTAL_SLIDES)

    cls       = ml.get("classification","—")
    sp        = ml.get("success_probability",0)*100
    fp        = ml.get("probabilities",{}).get("failure",0)*100
    up        = ml.get("probabilities",{}).get("uncertain",0)*100
    risk      = ml.get("risk_level","—")
    pred_fund = ml.get("predicted_funding_usd",0) or 0
    accent    = _ACCENT_MAP.get(cls, YELLOW)

    # Classification hero
    _rect(slide, 0.3, 0.7, 4.5, 1.15, BG_CARD)
    _rect(slide, 0.3, 0.7, 0.09, 1.15, accent)
    _txt(slide, "CLASSIFICATION", 0.5, 0.78, 4.2, 0.3, size=8, color=GRAY, bold=True)
    _txt(slide, cls, 0.5, 1.05, 4.2, 0.65, size=30, bold=True, color=accent)

    # Stat cards — top row
    cards_top = [
        ("Success Probability", f"{sp:.1f}%",    GREEN),
        ("Failure Probability", f"{fp:.1f}%",    RED),
        ("Uncertain",           f"{up:.1f}%",    YELLOW),
    ]
    for i, (lbl, val, col) in enumerate(cards_top):
        x = 5.1 + i * 1.65
        _stat_card(slide, x, 0.7, 1.55, 1.15, lbl, val, col)

    # Stat cards — bottom row
    _stat_card(slide, 0.3, 2.05, 4.5, 1.2, "Risk Level", risk,
               RED if risk=="High" else GREEN)
    _stat_card(slide, 5.1, 2.05, 4.7, 1.2,
               "Predicted Next Funding",
               f"${pred_fund/1e6:.2f}M" if pred_fund else "N/A", PURPLE)

    # Probability bar chart (visual)
    _txt(slide, "Probability Breakdown", 0.3, 3.45, 5, 0.3,
         size=10, bold=True, color=WHITE)
    bar_y = 3.82
    for label, pct, col in [("Success", sp, GREEN), ("Failure", fp, RED), ("Uncertain", up, YELLOW)]:
        bar_w = max(0.1, (pct / 100) * 8.8)
        _rect(slide, 0.3, bar_y, 8.8, 0.3, BG_CARD)
        _rect(slide, 0.3, bar_y, bar_w, 0.3, col)
        _txt(slide, f"{label}  {pct:.1f}%", 0.35, bar_y+0.04, 4, 0.22,
             size=8, color=WHITE, bold=True)
        bar_y += 0.42


# ── Slide 4 · Risk Factors ────────────────────────────────────
def _slide_risks(prs, risks, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "3. Risk Factors", n, TOTAL_SLIDES)

    if not risks:
        _txt(slide, "No specific risk factors identified.", 0.5, 2.5, 9, 0.5,
             size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
        return

    # Two-column layout for up to 8 risks
    left  = risks[:4]
    right = risks[4:8]

    for col_i, col_risks in enumerate([left, right]):
        x_base = 0.3 if col_i == 0 else 5.1
        for i, risk in enumerate(col_risks):
            y = 0.72 + i * 1.1
            _rect(slide, x_base, y, 4.6, 0.95, BG_CARD)
            _rect(slide, x_base, y, 0.07, 0.95, RED)
            _txt(slide, f"⚠  Risk {(col_i*4)+i+1}", x_base+0.18, y+0.06, 4.2, 0.28,
                 size=8, color=RED, bold=True)
            _txt(slide, _trunc(risk, 90), x_base+0.18, y+0.33, 4.25, 0.55,
                 size=9, color=GRAY_L, wrap=True)


# ── Slide 5 · Competitor Analysis ────────────────────────────
def _slide_competitors(prs, competitors, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "4. Competitor Analysis", n, TOTAL_SLIDES)

    if not competitors:
        _txt(slide, "No competitor data available.", 0.5, 2.5, 9, 0.5,
             size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
        return

    for i, c in enumerate(competitors[:5]):
        doc_text = c.get("document", c.get("name", c.get("summary", str(c))))
        meta     = c.get("metadata", {})
        industry = meta.get("industry", "")
        funding  = meta.get("funding", 0)

        y = 0.68 + i * 0.94
        _rect(slide, 0.3, y, 9.4, 0.82, BG_CARD)
        _rect(slide, 0.3, y, 0.07, 0.82, PURPLE)

        # Index circle
        _rect(slide, 0.45, y+0.17, 0.42, 0.42, PURPLE)
        _txt(slide, str(i+1), 0.45, y+0.17, 0.42, 0.42,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Text content
        _txt(slide, _trunc(str(doc_text), 110), 1.0, y+0.06, 6.0, 0.45,
             size=9, color=GRAY_L, wrap=True)

        # Right badges
        if industry:
            _rect(slide, 7.2, y+0.1, 1.3, 0.3, PURPLE_D)
            _txt(slide, industry[:18], 7.23, y+0.13, 1.24, 0.24,
                 size=7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        if funding:
            _rect(slide, 8.65, y+0.1, 0.9, 0.3, BG_MID)
            _txt(slide, f"${funding/1e6:.1f}M", 8.67, y+0.13, 0.86, 0.24,
                 size=7, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

        # Sub-line
        sub = _trunc(str(doc_text), 200)[110:]
        if sub:
            _txt(slide, sub, 1.0, y+0.5, 6.2, 0.28, size=8, color=GRAY, wrap=True)


# ── Slide 6 · Team Hierarchy ──────────────────────────────────
def _slide_hierarchy(prs, hierarchy, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "5. Team Hierarchy & Org Structure", n, TOTAL_SLIDES)

    ceo   = hierarchy.get("ceo_title", "CEO")
    total = hierarchy.get("total_employees", "—")
    depts = hierarchy.get("departments", [])
    gaps  = hierarchy.get("hiring_gaps", [])

    # CEO banner
    _rect(slide, 0.3, 0.65, 9.4, 0.52, PURPLE)
    _txt(slide, f"CEO / Founder: {ceo}   •   Total Employees: {total}",
         0.5, 0.75, 9.0, 0.34, size=12, bold=True, color=WHITE)

    # Departments — up to 5 in two columns
    left_d  = depts[:3]
    right_d = depts[3:6]

    for col_i, col_depts in enumerate([left_d, right_d]):
        x_base = 0.3 if col_i == 0 else 5.05
        for i, d in enumerate(col_depts):
            y = 1.3 + i * 1.28
            _rect(slide, x_base, y, 4.6, 1.15, BG_CARD)
            _rect(slide, x_base, y, 0.07, 1.15, GREEN)

            name = d.get("name","")
            hc   = d.get("headcount","")
            roles = ", ".join(
    r["title"] if isinstance(r, dict) else str(r)
    for r in d.get("roles", [])[:2]
)
            skills= ", ".join(d.get("skills_needed",[])[:3])

            _txt(slide, f"{name}  ({hc} people)",
                 x_base+0.18, y+0.08, 4.2, 0.32, size=11, bold=True, color=WHITE)
            _txt(slide, f"Roles: {roles}",
                 x_base+0.18, y+0.42, 4.2, 0.28, size=8, color=GRAY_L)
            _txt(slide, f"Skills: {skills}",
                 x_base+0.18, y+0.72, 4.2, 0.28, size=8, color=GRAY, italic=True)

    # Hiring gaps bar
    if gaps:
        _rect(slide, 0.3, 5.12, 9.4, 0.35, BG_MID)
        _txt(slide, "Gaps: " + "  •  ".join(gaps[:5]),
             0.45, 5.17, 9.0, 0.25, size=9, color=ORANGE, bold=True)


# ── Slide 7 · Recommended Next Hires ─────────────────────────
def _slide_next_hires(prs, hierarchy, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "5b. Recommended Next Hires", n, TOTAL_SLIDES)

    next_hires = hierarchy.get("recommended_next_hires", [])
    insight    = hierarchy.get("org_insight", "")

    if not next_hires:
        _txt(slide, "No hiring recommendations available.", 0.5, 2.5, 9, 0.5,
             size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
        return

    for i, h in enumerate(next_hires[:4]):
        role     = h.get("role","")
        priority = h.get("priority","")
        reason   = h.get("reason","")
        p_color  = RED if priority=="High" else (YELLOW if priority=="Medium" else GREEN)

        y = 0.68 + i * 1.1
        _rect(slide, 0.3, y, 9.4, 0.95, BG_CARD)
        _rect(slide, 0.3, y, 0.07, 0.95, p_color)

        _txt(slide, role, 0.5, y+0.08, 5.5, 0.38, size=14, bold=True, color=WHITE)

        # Priority pill
        _rect(slide, 6.2, y+0.1, 1.3, 0.32, p_color)
        _txt(slide, priority, 6.22, y+0.13, 1.26, 0.26,
             size=9, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

        _txt(slide, _trunc(reason, 100), 0.5, y+0.52, 8.8, 0.35,
             size=9, color=GRAY_L, italic=True, wrap=True)

    # Org insight
    if insight:
        _rect(slide, 0.3, H-0.85, 9.4, 0.72, BG_MID)
        _txt(slide, f"💡  {_trunc(insight, 160)}",
             0.45, H-0.78, 9.1, 0.55, size=9, color=GRAY_L, italic=True, wrap=True)


# ── Slide 8 · Hiring Guide Sequence + Profiles ───────────────
def _slide_hiring_guide(prs, hiring_guide, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "6. Hiring Guide & Role Profiles", n, TOTAL_SLIDES)

    seq      = hiring_guide.get("hiring_sequence", [])
    profiles = hiring_guide.get("hiring_profiles", [])
    culture  = hiring_guide.get("culture_fit_signals", [])

    # LEFT — hiring sequence
    _rect(slide, 0.3, 0.65, 4.5, 0.4, BG_MID)
    _txt(slide, "HIRING SEQUENCE", 0.4, 0.72, 4.3, 0.28, size=9, bold=True, color=PURPLE)

    for i, s in enumerate(seq[:5]):
        y = 1.12 + i * 0.78
        _rect(slide, 0.3, y, 4.5, 0.68, BG_CARD)
        _rect(slide, 0.3, y, 0.07, 0.68, PURPLE)
        _txt(slide, f"#{s.get('order',i+1)}  {s.get('role','')}",
             0.48, y+0.06, 4.1, 0.3, size=10, bold=True, color=WHITE)
        _txt(slide, _trunc(s.get("rationale",""), 70),
             0.48, y+0.38, 4.1, 0.26, size=8, color=GRAY, wrap=True)

    # RIGHT — top 3 hiring profiles
    _rect(slide, 5.0, 0.65, 4.7, 0.4, BG_MID)
    _txt(slide, "ROLE PROFILES", 5.1, 0.72, 4.5, 0.28, size=9, bold=True, color=GREEN)

    for i, p in enumerate(profiles[:4]):
        role     = p.get("role","")
        priority = p.get("priority","")
        dept     = p.get("department","")
        salary   = p.get("salary_range","")
        p_color  = RED if priority=="High" else (YELLOW if priority=="Medium" else GREEN)
        y = 1.12 + i * 1.05

        _rect(slide, 5.0, y, 4.7, 0.92, BG_CARD)
        _rect(slide, 5.0, y, 0.07, 0.92, p_color)
        _txt(slide, role, 5.18, y+0.05, 3.5, 0.32, size=11, bold=True, color=WHITE)
        _txt(slide, f"{dept}  •  {salary}", 5.18, y+0.4, 3.8, 0.24, size=8, color=GRAY)

        must = p.get("must_have_skills", [])[:3]
        if must:
            _txt(slide, "Skills: " + " · ".join(must),
                 5.18, y+0.62, 3.8, 0.24, size=8, color=GRAY_L, italic=True)


# ── Slide 9 · Hiring Detail — Skills, Responsibilities ───────
def _slide_hiring_detail(prs, hiring_guide, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "6b. Hiring Detail & Culture", n, TOTAL_SLIDES)

    profiles  = hiring_guide.get("hiring_profiles", [])
    culture   = hiring_guide.get("culture_fit_signals", [])
    onboarding= hiring_guide.get("onboarding_tips", "")

    # Show detail for first 2 profiles side by side
    for col_i, p in enumerate(profiles[:2]):
        x = 0.3 if col_i == 0 else 5.05
        role  = p.get("role","")
        why   = p.get("why_critical","")
        must  = p.get("must_have_skills",[])
        nice  = p.get("nice_to_have_skills",[])
        resp  = p.get("key_responsibilities",[])
        sigs  = p.get("interview_signals",[])
        seniority = p.get("seniority","")
        exp   = p.get("experience_years","")

        _rect(slide, x, 0.65, 4.6, 3.95, BG_CARD)
        _rect(slide, x, 0.65, 0.07, 3.95, GREEN)

        _txt(slide, role, x+0.18, 0.72, 4.25, 0.38, size=13, bold=True, color=WHITE)
        _txt(slide, f"{seniority}  •  {exp} yrs",
             x+0.18, 1.1, 4.25, 0.28, size=8, color=GRAY)

        if why:
            _txt(slide, _trunc(why, 90), x+0.18, 1.4, 4.25, 0.38,
                 size=8, color=GRAY_L, italic=True, wrap=True)

        y_cur = 1.82
        if must:
            _txt(slide, "Must-Have:", x+0.18, y_cur, 4.25, 0.25, size=8, bold=True, color=GREEN)
            _txt(slide, " · ".join(must[:4]), x+0.18, y_cur+0.25, 4.25, 0.3,
                 size=8, color=GRAY_L, wrap=True)
            y_cur += 0.6
        if nice:
            _txt(slide, "Nice to Have:", x+0.18, y_cur, 4.25, 0.25, size=8, bold=True, color=YELLOW)
            _txt(slide, " · ".join(nice[:3]), x+0.18, y_cur+0.25, 4.25, 0.3,
                 size=8, color=GRAY, wrap=True)
            y_cur += 0.6
        if resp:
            _txt(slide, "Responsibilities:", x+0.18, y_cur, 4.25, 0.25, size=8, bold=True, color=PURPLE)
            for r in resp[:2]:
                _txt(slide, f"→ {_trunc(r, 60)}", x+0.18, y_cur+0.25, 4.25, 0.28,
                     size=8, color=GRAY_L, wrap=True)
                y_cur += 0.28
        if sigs:
            _txt(slide, "Interview Signals: " + " | ".join(sigs[:2]),
                 x+0.18, 4.2, 4.25, 0.28, size=8, color=ORANGE, italic=True)

    # Culture + Onboarding strip at bottom
    _rect(slide, 0.3, 4.72, 9.4, 0.72, BG_MID)
    if culture:
        culture_text = "  ★  ".join(culture[:3])
        _txt(slide, f"Culture: {culture_text}", 0.45, 4.78, 9.0, 0.3,
             size=8, bold=True, color=PURPLE)
    if onboarding:
        _txt(slide, f"Onboarding: {_trunc(onboarding, 130)}",
             0.45, 5.05, 9.0, 0.28, size=8, color=GRAY_L, italic=True)


# ── Slide 10 · Strategic Analysis & 30-Day Plan ──────────────
def _slide_analysis(prs, analysis, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _slide_header(slide, "7. Strategic Analysis & 30-Day Action Plan", n, TOTAL_SLIDES)

    if not analysis:
        _txt(slide, "Strategic analysis not available.", 0.5, 2.5, 9, 0.5,
             size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)
        return

    # Parse action items and key points from analysis text
    bullets = []
    headers = []
    for line in analysis.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith(("- ", "* ", "• ")):
            bullets.append(_clean(line[2:]))
        elif line.startswith(("##", "**")):
            headers.append(_clean(line.replace("*","").replace("#","").strip()))

    # Left panel — key strategic themes
    _rect(slide, 0.3, 0.65, 4.5, 4.75, BG_CARD)
    _txt(slide, "STRATEGIC THEMES", 0.45, 0.72, 4.2, 0.3,
         size=9, bold=True, color=PURPLE)
    for i, h in enumerate(headers[:6]):
        y = 1.1 + i * 0.62
        _rect(slide, 0.35, y, 4.4, 0.52, BG_MID)
        _rect(slide, 0.35, y, 0.06, 0.52, PURPLE)
        _txt(slide, _trunc(h, 65), 0.52, y+0.09, 4.1, 0.35,
             size=9, bold=True, color=WHITE, wrap=True)

    # Right panel — action bullets
    _rect(slide, 5.0, 0.65, 4.7, 4.75, BG_CARD)
    _txt(slide, "ACTION ITEMS", 5.15, 0.72, 4.4, 0.3,
         size=9, bold=True, color=GREEN)
    for i, b in enumerate(bullets[:8]):
        y = 1.1 + i * 0.54
        _txt(slide, f"→  {_trunc(b, 72)}", 5.15, y, 4.45, 0.46,
             size=9, color=GRAY_L, wrap=True)

    # If no headers/bullets found, dump raw text in two chunks
    if not headers and not bullets:
        lines = [_clean(l) for l in analysis.split("\n") if l.strip()]
        _multiline(slide, [_trunc(l, 80) for l in lines[:10]],
                   0.35, 0.72, 4.4, 4.5, size=9, color=GRAY_L)
        _multiline(slide, [_trunc(l, 80) for l in lines[10:20]],
                   5.05, 0.72, 4.55, 4.5, size=9, color=GRAY_L)


# ── Slide 11 · Closing ────────────────────────────────────────
def _slide_closing(prs, ui, ml):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    cls       = ml.get("classification","Uncertain")
    sp        = ml.get("success_probability",0)*100
    pred_fund = ml.get("predicted_funding_usd",0) or 0
    risk      = ml.get("risk_level","—")
    accent    = _ACCENT_MAP.get(cls, YELLOW)

    # Left decoration
    _rect(slide, 0, 0, 0.12, H, PURPLE)

    _txt(slide, "Thank You", 0.4, 0.6, 7, 0.7,
         size=36, bold=True, color=WHITE)
    _txt(slide, "BizGenius Startup Intelligence", 0.4, 1.35, 6, 0.45,
         size=16, color=GRAY, italic=True)
    _txt(slide, ui.get("domain",""), 0.4, 1.85, 5, 0.38,
         size=14, color=PURPLE, bold=True)

    # Summary strip
    _rect(slide, 0.4, 2.5, 9.3, 1.65, BG_CARD)
    _txt(slide, "SUMMARY", 0.6, 2.6, 9, 0.32, size=9, bold=True, color=PURPLE)
    summary_items = [
        (f"Classification: {cls}", accent),
        (f"Success Probability: {sp:.1f}%", GREEN),
        (f"Risk Level: {risk}", RED if risk=="High" else GREEN),
        (f"Predicted Funding: ${pred_fund/1e6:.2f}M", PURPLE),
    ]
    for i, (text, col) in enumerate(summary_items):
        x = 0.6 + i * 2.35
        _txt(slide, text, x, 2.97, 2.25, 0.9, size=9, color=col, bold=True, wrap=True)

    _txt(slide, f"Generated by BizGenius  •  {datetime.now().strftime('%B %d, %Y')}  •  Confidential",
         0.4, H-0.45, 9.2, 0.35, size=9, color=GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# PUBLIC ENTRY POINT
# ════════════════════════════════════════════════════════════════

def generate_pitch_deck(data: dict) -> str:
    """
    Build the full pitch deck and return path to the .pptx file.
    Caller must delete the temp file after streaming.
    """
    ui          = data["user_input"]
    ml          = data["ml_results"]
    analysis    = data.get("analysis", "")
    competitors = data.get("competitors", [])
    hierarchy   = data.get("hierarchy", {})
    risks       = data.get("probable_risks", [])
    hiring_guide= data.get("hiring_guide", {})

    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    # Build all slides in report order
    _slide_cover(prs, ui, ml)                        # Cover
    _slide_snapshot(prs, ui, 2)                      # 1. Startup Snapshot
    _slide_ml(prs, ml, 3)                            # 2. ML Prediction
    _slide_risks(prs, risks, 4)                      # 3. Risk Factors
    _slide_competitors(prs, competitors, 5)          # 4. Competitors
    _slide_hierarchy(prs, hierarchy, 6)              # 5. Team Hierarchy
    _slide_next_hires(prs, hierarchy, 7)             # 5b. Next Hires
    if hiring_guide:
        _slide_hiring_guide(prs, hiring_guide, 8)   # 6. Hiring Guide
        _slide_hiring_detail(prs, hiring_guide, 9)  # 6b. Hiring Detail
    _slide_analysis(prs, analysis, 10)               # 7. Strategic Analysis
    _slide_closing(prs, ui, ml)                      # Closing

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return tmp.name